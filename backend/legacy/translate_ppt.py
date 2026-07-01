import os
import gc 
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
# ---------------------- LangChain 依赖 ----------------------
from langchain_core.tools import ToolException
from concurrent.futures import ThreadPoolExecutor, as_completed
from translate_text import translate_text, _translation_cache
from dotenv import load_dotenv
load_dotenv()
import os
from pathlib import Path

# 获取当前脚本所在目录的父级目录（即 SIS_Agent 根目录）
SIS_AGENT_ROOT = Path(__file__).parent.parent

# ---------------------- 全局配置（你的默认目录） ----------------------
DEFAULT_INPUT_DIR = os.getenv("TRANSLATE_INPUT_DIR") or str(SIS_AGENT_ROOT / "translate" / "input")
DEFAULT_OUTPUT_DIR = os.getenv("TRANSLATE_OUTPUT_DIR") or str(SIS_AGENT_ROOT / "translate" / "output")
DEFAULT_TARGET_LANG = os.getenv("TRANSLATE_TARGET_LANG") or "日语"
DEFAULT_DELAY = float(os.getenv("TRANSLATE_DELAY") or 1.2)
MAX_WORKERS = int(os.getenv("MAX_WORKERS") or 6)


def _collect_translatable(shapes, bucket):
    """
    递归收集需要翻译的单元（文本框、表格单元格），并进入组合形状(GroupShape)。
    修复要点：原代码 `for shape in slide.shapes` 不会进入组合形状，
             组合内的文字会整批漏翻。这里递归处理组合。
    bucket: list，收集到的形状/单元格对象。
    """
    for shape in shapes:
        # 组合形状：递归进入
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_translatable(shape.shapes, bucket)
            continue
        # 表格
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        bucket.append(cell)
            continue
        # 图表内文本（分类/系列名等）暂不在此处理——如需可扩展 shape.chart
        # 普通文本框 / 占位符
        if hasattr(shape, "text") and shape.text.strip():
            bucket.append(shape)


def _apply_jp_font(tf, jp_font):
    """给 text_frame 设置日语字体，保证显示。失败不影响主流程。"""
    try:
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass  # 部分容器（如表格单元格）可能不支持 auto_size
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.name = jp_font
                if run.font.size is None:
                    run.font.size = Pt(10)
    except Exception:
        pass


def translate_ppt_file(file_name: str, workspace_dir: str = DEFAULT_INPUT_DIR, 
                       target_lang: str = DEFAULT_TARGET_LANG, delay: float = DEFAULT_DELAY) -> str:
    """
    翻译PPT文件（并行版）：按幻灯片分组并行翻译，每组内顺序翻译，保持上下文连贯。
    """

    input_path = os.path.abspath(os.path.join(workspace_dir, file_name))
    output_path = os.path.abspath(os.path.join(workspace_dir, f"{os.path.splitext(file_name)[0]}_{target_lang}{os.path.splitext(file_name)[1]}"))
    
    # 检查文件
    if not os.path.exists(input_path):
        raise ToolException(f"文件不存在: {input_path}")

    prs = Presentation(input_path)
    total_slides = len(prs.slides)
    
    # 获取系统可用日语字体（用于后续设置）
    jp_fonts = ["Meiryo UI", "MS Gothic", "MS Mincho", "SimSun", "Arial Unicode MS"]
    jp_font = jp_fonts[0]
    for font in jp_fonts:
        try:
            from pptx.util import Font
            Font(name=font)
            jp_font = font
            break
        except:
            continue
    
    # 收集每个幻灯片的数据：幻灯片对象、需要翻译的形状列表、上下文
    slides_data = []
    for i, slide in enumerate(prs.slides):
        ctx = {
            "file_name": file_name,
            "slide_num": i + 1,
            "total_slides": total_slides,
            "translated_segments": [],
            "term_requirements": "TBOX译为TBOX、TSU译为TSU、CAN总线译为CANバス（日语）/CAN Bus（英语），公司名称不翻译"
        }
        shapes_to_translate = []
        # 递归收集（含组合形状、表格）
        _collect_translatable(slide.shapes, shapes_to_translate)
        slides_data.append((slide, shapes_to_translate, ctx))
    
    # 定义处理单个幻灯片的函数
    def process_slide(slide, shapes, ctx, target_lang, delay, jp_font):
        # 顺序翻译该幻灯片中的所有形状/单元格
        for shape in shapes:
            original_text = shape.text
            # 修复：逐单元兜底——单条翻译失败时保留原文，不再让整体中断
            try:
                translated = translate_text(original_text, target_lang, delay, context=ctx)
            except Exception as e:
                print(f"[跳过] 第{ctx['slide_num']}页某单元翻译失败，保留原文: {e}")
                translated = original_text
            if translated is None:
                translated = original_text
            shape.text = translated
            # 设置字体（保证显示）；文本框与表格单元格均有 text_frame
            if hasattr(shape, "text_frame"):
                _apply_jp_font(shape.text_frame, jp_font)
        return True
    
    # 并行处理每个幻灯片（不同幻灯片是不同的 part，互不干扰）
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = []
        for slide, shapes, ctx in slides_data:
            future = executor.submit(process_slide, slide, shapes, ctx, target_lang, delay, jp_font)
            futures.append(future)
        # 等待所有任务完成
        for future in as_completed(futures):
            future.result()
    
    # 保存最终文件
    prs.save(output_path)
    
    # 清理
    # _translation_cache.clear()
    gc.collect()
    
    return f"PPT翻译完成！输出路径: {output_path}"